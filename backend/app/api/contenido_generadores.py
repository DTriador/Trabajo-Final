# backend/app/api/contenido_generadores.py
"""
Lógica de generación de cada tipo de documento (llamada a la IA + armado
del .docx / .pptx en memoria). Extraído sin cambios de lógica desde
router_contenido.py: cada función de acá es el "cuerpo" que antes vivía
pegado dentro del endpoint correspondiente. Los endpoints en
router_contenido.py ahora solo parsean el Form, llaman a la función acá
definida, y suben el archivo con process_and_upload (o, en el caso de
presentación, arman la respuesta final).
"""
import os
import re
import json
import uuid
import math
from io import BytesIO
from typing import Optional
from datetime import datetime

from fastapi import HTTPException, UploadFile
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.core.database import supabase
from app.services.ai_service import SYSTEM_PROMPT_PPTX
from app.services.rag_orchestrator import RAGOrchestrator
from app.api.contenido_helpers import _datos_escuela_materia, _contexto_biblio, _encabezado_documento
from app.utils.visual_theme import THEME

PPT_SLIDE_WIDTH = Inches(13.33)
PPT_SLIDE_HEIGHT = Inches(7.5)
PPT_MARGIN_LEFT = Inches(0.65)
PPT_MARGIN_RIGHT = Inches(0.65)
PPT_MARGIN_TOP = Inches(0.6)
PPT_MARGIN_BOTTOM = Inches(0.45)
PPT_TITLE_FONT_SIZE = 24
PPT_SUBTITLE_FONT_SIZE = 18
PPT_BODY_FONT_SIZE = 18
PPT_MIN_BODY_FONT_SIZE = 14
PPT_MAX_BULLETS_PER_SLIDE = 6


def _coerce_ppt_text(value, fallback: str = "") -> str:
    if value is None:
        return fallback
    text = str(value).strip()
    text = re.sub(r"\s+", " ", text)
    return text or fallback


def _clean_ppt_items(items):
    if items is None:
        return []
    if isinstance(items, str):
        items = [items]
    if not isinstance(items, list):
        return []

    limpio = []
    vistos = set()
    for item in items:
        if item is None:
            continue
        texto = _coerce_ppt_text(item)
        if not texto:
            continue
        clave = texto.lower()
        if clave in vistos:
            continue
        vistos.add(clave)
        limpio.append(texto)
    return limpio


def _split_ppt_items(items, max_per_slide: int = PPT_MAX_BULLETS_PER_SLIDE):
    items = _clean_ppt_items(items)
    if not items:
        return [[]]
    return [items[i:i + max_per_slide] for i in range(0, len(items), max_per_slide)]


def _normalize_title(text: str, fallback: str = "Diapositiva") -> str:
    title = _coerce_ppt_text(text, fallback)
    if len(title) > 80:
        title = title[:77].rstrip() + "..."
    return title


def _normalize_ppt_slide_payload(slides):
    if not isinstance(slides, list):
        slides = []

    normalizadas = []
    for idx, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            continue

        titulo = _normalize_title(
            slide.get("titulo") or slide.get("subtitulo") or slide.get("title") or slide.get("titulo_slide") or f"Diapositiva {idx}",
            fallback=f"Diapositiva {idx}",
        )
        subtitulo = _coerce_ppt_text(slide.get("subtitle") or slide.get("subtitulo_secundario") or "")
        contenido = slide.get("contenido") or slide.get("bullets") or slide.get("items") or slide.get("content") or []
        items = _clean_ppt_items(contenido)

        if not items and slide.get("descripcion"):
            items = _clean_ppt_items([slide.get("descripcion")])

        if not items:
            normalizadas.append({"titulo": titulo, "subtitulo": subtitulo, "contenido": []})
            continue

        bloques = _split_ppt_items(items, max_per_slide=PPT_MAX_BULLETS_PER_SLIDE)
        for parte_idx, bloque in enumerate(bloques, start=1):
            slide_title = titulo if len(bloques) == 1 else f"{titulo} — Parte {parte_idx}"
            slide_subtitle = subtitulo if parte_idx == 1 else ""
            normalizadas.append({"titulo": slide_title, "subtitulo": slide_subtitle, "contenido": bloque})

    if not normalizadas:
        normalizadas = [{
            "titulo": "Introducción",
            "subtitulo": "",
            "contenido": ["Tema principal", "Objetivos clave", "Conclusiones relevantes"],
        }]

    return normalizadas


def _validate_slide_bounds(left, top, width, height, slide_width, slide_height, tolerance=0.2):
    return (
        left >= -tolerance and
        top >= -tolerance and
        left + width <= slide_width + tolerance and
        top + height <= slide_height + tolerance
    )


def _estimate_text_block_height(text: str, font_size: int = PPT_BODY_FONT_SIZE):
    if not text:
        return Inches(0.25)
    estimated_chars_per_line = 38
    line_count = max(1, math.ceil(len(text) / estimated_chars_per_line))
    return Inches(0.25 + line_count * (font_size / 72) * 1.3)


def _add_textbox_text(box, text: str, font_size: int, bold: bool = False, align=PP_ALIGN.LEFT):
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    para = tf.paragraphs[0]
    para.text = text
    para.alignment = align
    if para.runs:
        para.runs[0].font.size = Pt(font_size)
        para.runs[0].font.bold = bold


def _create_slide_layout(prs, slide_data, index: int = 1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = _normalize_title(slide_data.get("titulo") or f"Diapositiva {index}", fallback=f"Diapositiva {index}")
    subtitle = _coerce_ppt_text(slide_data.get("subtitulo") or "")
    items = _clean_ppt_items(slide_data.get("contenido") or [])

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    content_left = PPT_MARGIN_LEFT
    content_right = slide_width - PPT_MARGIN_RIGHT
    content_width = content_right - content_left

    title_box = slide.shapes.add_textbox(PPT_MARGIN_LEFT, PPT_MARGIN_TOP, content_width, Inches(0.5))
    _add_textbox_text(title_box, title, PPT_TITLE_FONT_SIZE, bold=True)
    title_height = title_box.height

    body_top = PPT_MARGIN_TOP + title_height + Inches(0.15)
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(PPT_MARGIN_LEFT, body_top, content_width, Inches(0.35))
        _add_textbox_text(subtitle_box, subtitle, PPT_SUBTITLE_FONT_SIZE)
        body_top = subtitle_box.top + subtitle_box.height + Inches(0.12)

    if not items:
        return slide

    estimated_body = len(items) * Inches(0.34) + Inches(0.2)
    max_body_height = slide_height - body_top - PPT_MARGIN_BOTTOM
    if estimated_body > max_body_height:
        item_count = max(1, int(max_body_height / Inches(0.34)))
        items = items[:item_count]

    body_box = slide.shapes.add_textbox(PPT_MARGIN_LEFT, body_top, content_width, max_body_height)
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    body_tf.margin_left = Inches(0.1)
    body_tf.margin_right = 0
    body_tf.margin_top = 0
    body_tf.margin_bottom = 0

    for item in items:
        if len(item) > 140:
            item = item[:137].rstrip() + "..."
        p = body_tf.paragraphs[0] if not body_tf.paragraphs[0].text else body_tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        if p.runs:
            p.runs[0].font.size = Pt(PPT_BODY_FONT_SIZE)
            p.runs[0].font.bold = False
        p.space_after = Pt(5)

    if not _validate_slide_bounds(PPT_MARGIN_LEFT, PPT_MARGIN_TOP, content_width, slide_height - PPT_MARGIN_TOP - PPT_MARGIN_BOTTOM, slide_width, slide_height):
        raise ValueError("La diapositiva generada excede los límites del layout.")

    return slide


# ── apunte ────────────────────────────────────────────────────────────────────

async def generar_apunte_docx(
    tema: str,
    id_docente: str,
    file: Optional[UploadFile],
    id_escuela: Optional[str],
    id_curso: Optional[str],
    fecha: Optional[str],
):
    """Genera el .docx de un apunte. Devuelve (bytes, nombre_archivo)."""
    from docx import Document

    nombre_escuela, nombre_materia, division, contenido_minimo, bibliografia = \
        _datos_escuela_materia(id_escuela, id_curso)

    SYSTEM_PROMPT_APUNTE = """
    Sos un asistente pedagógico experto en sintetizar material académico.
    A partir del contenido proporcionado, generá un APUNTE de estudio claro,
    enfocado en el tema indicado, listo para que un alumno lo use para repasar.

    Devolvé EXCLUSIVAMENTE un JSON válido con esta estructura:
    {
      "titulo": "Título del apunte",
      "introduccion": "Párrafo breve que contextualiza el tema",
      "secciones": [
        {
          "subtitulo": "Nombre de la sección",
          "contenido": "Explicación clara y desarrollada del concepto",
          "puntos_clave": ["idea 1", "idea 2", "idea 3"]
        }
      ],
      "glosario": [
        {"termino": "Concepto", "definicion": "Definición breve"}
      ],
      "conclusion": "Cierre con las ideas más importantes para recordar"
    }
    Mínimo 4 secciones y 5 términos en el glosario.
    """

    prompt_base = (
        f"Generá un APUNTE / RESUMEN de estudio sobre el tema '{tema}' "
        f"para la materia '{nombre_materia}' de '{nombre_escuela}'."
    )

    if file is not None:
        pdf_content = await file.read()
        datos_json = await RAGOrchestrator.get_context_from_file_and_generate(
            pdf_content, prompt_base, SYSTEM_PROMPT_APUNTE, id_docente=id_docente
        )
    else:
        ctx = _contexto_biblio(contenido_minimo, bibliografia)
        datos_json = await RAGOrchestrator.get_context_and_generate(
            f"{prompt_base}\n\n{ctx}", SYSTEM_PROMPT_APUNTE, id_docente=id_docente
        )

    doc = Document()
    _encabezado_documento(doc, "Apunte", nombre_materia, tema, nombre_escuela, fecha_str=fecha)

    if datos_json.get("introduccion"):
        doc.add_heading("Introducción", level=1)
        doc.add_paragraph(datos_json["introduccion"])

    for seccion in datos_json.get("secciones", []):
        doc.add_heading(seccion.get("subtitulo", "Sección"), level=1)
        if seccion.get("contenido"):
            doc.add_paragraph(seccion["contenido"])
        for punto in seccion.get("puntos_clave", []):
            doc.add_paragraph(punto, style="List Bullet")

    if datos_json.get("glosario"):
        doc.add_heading("Glosario", level=1)
        for item in datos_json["glosario"]:
            p = doc.add_paragraph()
            p.add_run(f"{item.get('termino', '')}: ").bold = True
            p.add_run(item.get("definicion", ""))

    if datos_json.get("conclusion"):
        doc.add_heading("Conclusión", level=1)
        doc.add_paragraph(datos_json["conclusion"])

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    nombre_archivo = f"Apunte_{nombre_escuela}_{nombre_materia}".replace(" ", "_")
    return buffer.getvalue(), nombre_archivo


# ── preguntas guía ────────────────────────────────────────────────────────────

async def generar_preguntas_docx(
    tema: str,
    id_docente: str,
    nombre_guia: str,
    numero_preguntas: int,
    file: Optional[UploadFile],
    id_escuela: Optional[str],
    id_curso: Optional[str],
    fecha: Optional[str],
):
    """Genera el .docx de una guía de preguntas. Devuelve (bytes, nombre_archivo)."""
    from docx import Document

    if numero_preguntas < 1 or numero_preguntas > 25:
        raise HTTPException(status_code=400, detail="La cantidad de preguntas debe estar entre 1 y 25.")

    nombre_escuela, nombre_materia, division, contenido_minimo, bibliografia = \
        _datos_escuela_materia(id_escuela, id_curso)

    SYSTEM_PROMPT_PREGUNTAS = f"""
    Sos un asistente pedagógico experto en generar guías de preguntas
    a partir de un material de lectura.

    Generá EXACTAMENTE {numero_preguntas} preguntas guía sobre el tema indicado.
    Combiná preguntas de comprensión literal, de análisis y de reflexión.
    No inventes contenido que no esté en el material provisto.

    IMPORTANTE: Devolvé EXCLUSIVAMENTE un JSON válido. Estructura EXACTA:
    {{
      "titulo": "Título de la guía",
      "introduccion": "Indicaciones breves para el alumno",
      "preguntas": [
        {{"numero": 1, "pregunta": "Texto de la pregunta", "respuesta_sugerida": "Respuesta orientativa"}}
      ]
    }}
    El array 'preguntas' debe tener EXACTAMENTE {numero_preguntas} elementos.
    """

    prompt_base = (
        f"Generá una guía de {numero_preguntas} preguntas sobre el tema '{tema}' "
        f"para la materia '{nombre_materia}' de '{nombre_escuela}'."
    )

    if file is not None:
        pdf_content = await file.read()
        if not pdf_content:
            raise HTTPException(status_code=400, detail="El PDF está vacío.")
        datos_json = await RAGOrchestrator.get_context_from_file_and_generate(
            pdf_content, prompt_base, SYSTEM_PROMPT_PREGUNTAS, id_docente=id_docente
        )
    else:
        ctx = _contexto_biblio(contenido_minimo, bibliografia)
        datos_json = await RAGOrchestrator.get_context_and_generate(
            f"{prompt_base}\n\n{ctx}", SYSTEM_PROMPT_PREGUNTAS, id_docente=id_docente
        )

    # Normalización defensiva
    if isinstance(datos_json, str):
        txt = datos_json.strip()
        m = re.search(r"\{.*\}", txt, re.DOTALL)
        if m:
            try:
                datos_json = json.loads(m.group(0))
            except Exception:
                datos_json = {}
    if not isinstance(datos_json, dict):
        datos_json = {}

    preguntas    = datos_json.get("preguntas") or datos_json.get("questions") or datos_json.get("items") or []
    titulo       = datos_json.get("titulo") or datos_json.get("title") or nombre_guia
    introduccion = datos_json.get("introduccion") or datos_json.get("introduction") or ""

    if not preguntas:
        raise HTTPException(
            status_code=500,
            detail=f"El LLM no devolvió preguntas. Respuesta cruda: {str(datos_json)[:300]}"
        )

    doc = Document()
    _encabezado_documento(doc, "Guía de Preguntas", nombre_materia, tema, nombre_escuela, fecha_str=fecha)

    if introduccion:
        doc.add_heading("Indicaciones", level=1)
        doc.add_paragraph(introduccion)

    doc.add_heading("Preguntas", level=1)
    for i, item in enumerate(preguntas, start=1):
        if isinstance(item, dict):
            texto  = item.get("pregunta") or item.get("question") or item.get("texto") or ""
            numero = item.get("numero", i)
        else:
            texto, numero = str(item), i
        p = doc.add_paragraph()
        p.add_run(f"{numero}. ").bold = True
        p.add_run(texto)

    doc.add_page_break()
    doc.add_heading("Respuestas sugeridas (uso docente)", level=1)
    for i, item in enumerate(preguntas, start=1):
        if isinstance(item, dict):
            texto  = item.get("pregunta") or item.get("question") or ""
            resp   = item.get("respuesta_sugerida") or item.get("respuesta") or item.get("answer") or ""
            numero = item.get("numero", i)
        else:
            texto, resp, numero = str(item), "", i
        p = doc.add_paragraph()
        p.add_run(f"{numero}. {texto}\n").bold = True
        p.add_run(resp)

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    nombre_archivo = f"Guia_{nombre_escuela}_{nombre_materia}_{nombre_guia}".replace(" ", "_")
    return buffer.getvalue(), nombre_archivo


# ── examen ────────────────────────────────────────────────────────────────────

async def generar_examen_docx(
    id_docente: str,
    materia: Optional[str],
    fecha_examen: str,
    tipos: str,
    file: Optional[UploadFile],
    id_escuela: Optional[str],
    id_curso: Optional[str],
):
    """Genera el .docx de un examen. Devuelve (bytes, nombre_archivo, nombre_materia)."""
    from docx import Document

    try:
        tipos_dict = json.loads(tipos)
    except Exception:
        raise HTTPException(status_code=400, detail="Formato inválido en 'tipos'.")

    seleccionados = {
        k: v for k, v in tipos_dict.items()
        if v.get("activo") and int(v.get("cantidad", 0)) > 0
    }
    if not seleccionados:
        raise HTTPException(status_code=400, detail="Marcá al menos una actividad con cantidad > 0.")

    # Obtener escuela y materia desde la BD
    nombre_escuela, nombre_materia_db, division, contenido_minimo, bibliografia = \
        _datos_escuela_materia(id_escuela, id_curso)

    # El campo materia del Form tiene prioridad; si no viene, usar el de la BD
    nombre_materia = (materia or "").strip() or nombre_materia_db

    # Bibliografía de fallback cuando no hay PDF
    biblio_str = ""
    if file is None:
        biblio_str = _contexto_biblio(contenido_minimo, bibliografia)

    descripcion_tipos = "\n".join([f"- {k}: {v['cantidad']} ítems" for k, v in seleccionados.items()])

    SYSTEM_PROMPT_EXAMEN = f"""
    Sos un asistente pedagógico experto. Generá un examen escrito sobre la materia '{nombre_materia}'.
    Incluí EXACTAMENTE los siguientes tipos y cantidades:
    {descripcion_tipos}

    Devolvé EXCLUSIVAMENTE un JSON válido (sin markdown):
    {{
      "titulo": "Examen de {nombre_materia}",
      "consignas": [
        {{"tipo": "desarrollo|multiple|completar|verdadero_falso",
          "enunciado": "Texto de la consigna",
          "items": ["item1", "item2", ...]}}
      ]
    }}
    """

    if file is not None:
        prompt = f"Generá el examen de '{nombre_materia}' para la fecha {fecha_examen}, basado en el PDF adjunto."
        pdf_content = await file.read()
        if not pdf_content:
            raise HTTPException(status_code=400, detail="El PDF está vacío.")
        datos_json = await RAGOrchestrator.get_context_from_file_and_generate(
            pdf_content, prompt, SYSTEM_PROMPT_EXAMEN, id_docente=id_docente
        )
    else:
        prompt = (
            f"Generá el examen de '{nombre_materia}' para la fecha {fecha_examen}."
            + (f"\n\n{biblio_str}" if biblio_str else "")
        )
        datos_json = await RAGOrchestrator.get_context_and_generate(
            prompt,
            SYSTEM_PROMPT_EXAMEN,
            id_docente=id_docente,
        )

    if isinstance(datos_json, str):
        m = re.search(r"\{.*\}", datos_json, re.DOTALL)
        if m:
            try:
                datos_json = json.loads(m.group(0))
            except Exception:
                datos_json = {}
    if not isinstance(datos_json, dict):
        datos_json = {}

    consignas  = datos_json.get("consignas") or []
    titulo_doc = datos_json.get("titulo") or f"Examen de {nombre_materia}"

    if not consignas:
        raise HTTPException(status_code=500, detail=f"El LLM no devolvió consignas: {str(datos_json)[:300]}")

    doc = Document()
    _encabezado_documento(
        doc, "Examen", nombre_materia, nombre_materia,
        nombre_escuela, fecha_str=fecha_examen
    )

    doc.add_heading(titulo_doc, level=1)

    for idx, c in enumerate(consignas, start=1):
        tipo      = c.get("tipo", "")
        enunciado = c.get("enunciado", "")
        items     = c.get("items", []) or []
        doc.add_paragraph(f"{idx}) [{tipo}] {enunciado}", style="Heading 3")
        for j, it in enumerate(items, start=1):
            doc.add_paragraph(f"{j}. {it}")

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    nombre_archivo = f"Examen_{nombre_escuela}_{nombre_materia}".replace(" ", "_")
    return buffer.getvalue(), nombre_archivo, nombre_materia


# ── podcast ───────────────────────────────────────────────────────────────────

async def generar_podcast_docx(
    tema: str,
    id_docente: str,
    file: Optional[UploadFile],
    id_escuela: Optional[str],
    id_curso: Optional[str],
    fecha: Optional[str],
):
    """Genera el .docx de un guión de podcast. Devuelve (bytes, nombre_archivo)."""
    from docx import Document

    nombre_escuela, nombre_materia, division, contenido_minimo, bibliografia = \
        _datos_escuela_materia(id_escuela, id_curso)

    SYSTEM_PROMPT_PODCAST = """
    Sos un asistente pedagógico experto en crear guiones de podcast educativo.
    A partir del contenido proporcionado, generá un GUIÓN de podcast completo,
    entretenido y pedagógico, listo para grabar.

    Devolvé EXCLUSIVAMENTE un JSON válido con esta estructura:
    {
      "titulo": "Título del episodio",
      "duracion_estimada": "15-20 minutos",
      "segmentos": [
        {
          "tipo": "introduccion|desarrollo|reflexion|conclusion",
          "titulo": "Nombre del segmento",
          "guion": "Texto completo del segmento con indicaciones entre [corchetes]"
        }
      ]
    }
    Mínimo 4 segmentos. El guión debe ser natural, conversacional y pedagógico.
    """

    prompt_base = (
        f"Generá un guión de podcast educativo sobre el tema '{tema}' "
        f"para la materia '{nombre_materia}' de '{nombre_escuela}'."
    )

    if file is not None:
        pdf_content = await file.read()
        if not pdf_content:
            raise HTTPException(status_code=400, detail="El PDF está vacío.")
        datos_json = await RAGOrchestrator.get_context_from_file_and_generate(
            pdf_content, prompt_base, SYSTEM_PROMPT_PODCAST, id_docente=id_docente
        )
    else:
        ctx = _contexto_biblio(contenido_minimo, bibliografia)
        datos_json = await RAGOrchestrator.get_context_and_generate(
            f"{prompt_base}\n\n{ctx}", SYSTEM_PROMPT_PODCAST, id_docente=id_docente
        )

    # Normalización defensiva
    if isinstance(datos_json, str):
        m = re.search(r"\{.*\}", datos_json, re.DOTALL)
        if m:
            try:
                datos_json = json.loads(m.group(0))
            except Exception:
                datos_json = {}
    if not isinstance(datos_json, dict):
        datos_json = {}

    titulo_ep = datos_json.get("titulo") or f"Podcast: {tema}"
    duracion  = datos_json.get("duracion_estimada") or ""
    segmentos = datos_json.get("segmentos") or []

    if not segmentos:
        raise HTTPException(
            status_code=500,
            detail=f"El LLM no devolvió segmentos. Respuesta cruda: {str(datos_json)[:300]}"
        )

    doc = Document()
    _encabezado_documento(doc, "Podcast", nombre_materia, tema, nombre_escuela, fecha_str=fecha)

    doc.add_heading(titulo_ep, level=1)
    if duracion:
        p_dur = doc.add_paragraph()
        p_dur.add_run("Duración estimada: ").bold = True
        p_dur.add_run(duracion)
    doc.add_paragraph()

    for seg in segmentos:
        tipo_seg   = (seg.get("tipo") or "").upper()
        titulo_seg = seg.get("titulo") or tipo_seg
        guion_seg  = seg.get("guion") or ""
        doc.add_heading(f"[{tipo_seg}] {titulo_seg}", level=2)
        doc.add_paragraph(guion_seg)
        doc.add_paragraph()

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    nombre_archivo = f"Podcast_{nombre_escuela}_{nombre_materia}".replace(" ", "_")
    return buffer.getvalue(), nombre_archivo


# ── presentacion ──────────────────────────────────────────────────────────────

async def generar_presentacion_pptx(tema: str, id_docente: str, file: Optional[UploadFile]):
    """
    Genera un PPTX con contenido académico normalizado y layout controlado.
    El flujo actual llama a la IA para obtener estructura de diapositivas,
    limpia el contenido, particiona diapositivas largas y valida los límites del
    slide antes de guardar el archivo final.
    """
    contenido_pdf = ""
    if file:
        try:
            from pypdf import PdfReader
            pdf_bytes = await file.read()
            temp_path = f"/tmp/{uuid.uuid4()}.pdf"
            with open(temp_path, "wb") as f:
                f.write(pdf_bytes)
            reader = PdfReader(temp_path)
            contenido_pdf = "\n".join([p.extract_text() or "" for p in reader.pages])
            os.remove(temp_path)
        except Exception as e:
            print(f"⚠️ No pude leer el PDF: {e}")
            contenido_pdf = ""

    prompt_base = (
        f"Genera una presentación académica breve y visualmente clara sobre '{tema}'. "
        "Usa un estilo de diapositivas con título corto, subtítulo opcional y 3 a 6 ideas principales por slide. "
        "Prioriza contenido breve, conceptos resumidos y estructura fácil de leer."
    )

    if contenido_pdf:
        datos_json = await RAGOrchestrator.get_context_from_file_and_generate(
            file_content=contenido_pdf.encode("utf-8", errors="ignore"),
            user_prompt=prompt_base,
            system_instruction=SYSTEM_PROMPT_PPTX,
            id_docente=id_docente,
        )
    else:
        datos_json = await RAGOrchestrator.get_context_and_generate(
            user_prompt=prompt_base,
            system_instruction=SYSTEM_PROMPT_PPTX,
            id_docente=id_docente,
        )

    if isinstance(datos_json, str):
        texto = datos_json.strip()
        match = re.search(r"\{.*\}", texto, re.DOTALL)
        if match:
            try:
                datos_json = json.loads(match.group(0))
            except Exception:
                datos_json = {}
        else:
            try:
                datos_json = json.loads(texto)
            except Exception:
                datos_json = {}

    if not isinstance(datos_json, dict):
        datos_json = {}

    titulo_presentacion = _coerce_ppt_text(datos_json.get("titulo") or tema, tema)
    slides_raw = datos_json.get("slides") or []
    if not isinstance(slides_raw, list):
        slides_raw = []

    slides = _normalize_ppt_slide_payload(slides_raw)
    if not slides:
        slides = [{
            "titulo": titulo_presentacion,
            "subtitulo": "",
            "contenido": ["Objetivo principal", "Conceptos clave", "Conclusiones relevantes"],
        }]

    prs = Presentation()
    prs.slide_width = PPT_SLIDE_WIDTH
    prs.slide_height = PPT_SLIDE_HEIGHT

    titulo_slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = titulo_slide.background.fill
    background.solid()
    background.fore_color.rgb = 0xF5F7FF

    title_box = titulo_slide.shapes.add_textbox(PPT_MARGIN_LEFT, PPT_MARGIN_TOP, prs.slide_width - PPT_MARGIN_LEFT - PPT_MARGIN_RIGHT, Inches(0.7))
    _add_textbox_text(title_box, titulo_presentacion, PPT_TITLE_FONT_SIZE, bold=True)
    title_box.line.fill.background()
    title_box.fill.background()

    subtitle_box = titulo_slide.shapes.add_textbox(
        PPT_MARGIN_LEFT,
        title_box.top + title_box.height + Inches(0.15),
        prs.slide_width - PPT_MARGIN_LEFT - PPT_MARGIN_RIGHT,
        Inches(0.35),
    )
    _add_textbox_text(subtitle_box, f"Generado por Kōkua • {datetime.now().strftime('%d/%m/%Y')}", PPT_SUBTITLE_FONT_SIZE)
    subtitle_box.line.fill.background()
    subtitle_box.fill.background()

    for i, slide_data in enumerate(slides, start=1):
        _create_slide_layout(prs, slide_data, index=i)

    nombre_archivo = f"{titulo_presentacion.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    os.makedirs("archivos_generados", exist_ok=True)
    ruta_local = f"archivos_generados/{nombre_archivo}"
    prs.save(ruta_local)

    tamanio_mb = round(os.path.getsize(ruta_local) / (1024 * 1024), 3)

    try:
        with open(ruta_local, "rb") as f:
            supabase.storage.from_("documentos_docentes").upload(
                path=f"{id_docente}/{nombre_archivo}",
                file=f.read(),
                file_options={
                    "content-type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    "x-upsert": "true",
                }
            )
        signed = supabase.storage.from_("documentos_docentes").create_signed_url(
            f"{id_docente}/{nombre_archivo}", 60 * 60 * 24 * 7
        )
        url_publica = signed.get("signedURL") or signed.get("signedUrl") or signed.get("signed_url")
    except Exception as e:
        print(f"❌ No pude subir la presentación a Supabase Storage: {e}")
        raise HTTPException(
            status_code=500,
            detail=f"No se pudo guardar la presentación en el almacenamiento: {e}"
        )

    try:
        supabase.table("archivos_generados").insert({
            "id_docente":      id_docente,
            "nombre_archivo":  nombre_archivo,
            "tipo_formato":    "pptx",
            "sub_tipo":        "presentacion",
            "tema_especifico": tema,
            "prompt_origen":   f"Generar presentación sobre: {tema}",
            "fecha_creacion":  datetime.now().isoformat(),
            "categoria_ia":    "presentacion",
            "url_descarga":    url_publica,
            "uso_mb":          tamanio_mb,
            "descripcion":     f"Presentación generada automáticamente sobre el tema '{tema}'",
        }).execute()
    except Exception as e:
        print(f"⚠️ No pude registrar en BD: {e}")

    return {
        "status":         "success",
        "message":        "Presentación generada con éxito",
        "download_url":   url_publica,
        "nombre_archivo": nombre_archivo,
    }