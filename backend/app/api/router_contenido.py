# backend/app/api/router_contenido.py
from fastapi import APIRouter, HTTPException, UploadFile, File, Form
from typing import Optional
from app.utils.engines import FileEngine
from app.core.database import supabase
from app.services.rag_orchestrator import RAGOrchestrator
from app.services.ai_service import SYSTEM_PROMPT_XLSX
from app.api.generacion_utils import process_and_upload
import uuid, os, json, re
from io import BytesIO
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
import shutil
import subprocess
import tempfile
import httpx
from pathlib import Path
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from app.utils.storage import construir_ruta_storage

router = APIRouter()


# ── helpers internos ──────────────────────────────────────────────────────────

def _datos_escuela_materia(id_escuela: Optional[str], id_curso: Optional[str]):
    """
    Devuelve (nombre_escuela, nombre_materia, division, contenido_minimo, bibliografia).
    Consulta Supabase solo si se proveen los IDs.
    """
    nombre_escuela   = "Institución Educativa"
    nombre_materia   = "Materia General"
    division         = "-"
    contenido_minimo = ""
    bibliografia     = []

    if id_escuela:
        try:
            res = supabase.table("escuelas").select("nombre_escuela") \
                .eq("id_escuela", id_escuela).single().execute()
            nombre_escuela = (res.data or {}).get("nombre_escuela", nombre_escuela)
        except Exception as e:
            print(f"⚠️ No pude obtener escuela: {e}")

    if id_curso:
        try:
            res = supabase.table("cursos") \
                .select("nombre_materia,division,contenido_minimo,bibliografia") \
                .eq("id_curso", id_curso).single().execute()
            if res.data:
                nombre_materia   = res.data.get("nombre_materia",   nombre_materia)
                division         = res.data.get("division",         division)
                contenido_minimo = res.data.get("contenido_minimo", "") or ""
                bibliografia     = res.data.get("bibliografia",     []) or []
        except Exception as e:
            print(f"⚠️ No pude obtener curso: {e}")

    return nombre_escuela, nombre_materia, division, contenido_minimo, bibliografia


def _contexto_biblio(contenido_minimo: str, bibliografia: list) -> str:
    """Arma texto de contexto cuando no hay PDF."""
    partes = []
    if contenido_minimo:
        partes.append(f"Contenido mínimo de la materia:\n{contenido_minimo}")
    if bibliografia:
        items = [b if isinstance(b, str) else str(b) for b in bibliografia]
        partes.append("Bibliografía de la materia:\n" + "\n".join(f"- {i}" for i in items))
    return "\n\n".join(partes) if partes else \
        "Sin contenido específico. Generá basándote en el tema indicado."


def _encabezado_documento(doc, tipo_doc: str, nombre_materia: str, tema: str,
                           nombre_escuela: str, fecha_str: Optional[str] = None):
    """
    Agrega el encabezado institucional al documento Word.
    Formato (imagen de referencia):
        [Centrado, negrita, 18pt] {tipo_doc} de {nombre_materia}
        [Centrado, negrita, 14pt] Tema: {tema}
        <línea en blanco>
        Escuela: {nombre_escuela}
        Nombre del alumno: ______________________________
        Fecha: {fecha}
    """
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    # If fecha_str is None -> use current date; if it's an empty string -> keep it empty
    if fecha_str is None:
        fecha = datetime.now().strftime("%d/%m/%Y")
    else:
        fecha = fecha_str

    # Título principal centrado
    titulo_p = doc.add_paragraph()
    titulo_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_titulo = titulo_p.add_run(f"{tipo_doc} de {nombre_materia}")
    run_titulo.bold = True
    run_titulo.font.size = Pt(18)

    # Subtítulo Tema centrado
    tema_p = doc.add_paragraph()
    tema_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run_tema = tema_p.add_run(f"Tema: {tema}")
    run_tema.bold = True
    run_tema.font.size = Pt(14)

    doc.add_paragraph()  # espacio en blanco

    # Escuela
    p_esc = doc.add_paragraph()
    p_esc.add_run("Escuela: ").bold = True
    p_esc.add_run(nombre_escuela)

    # Nombre del alumno (siempre en blanco para que el alumno complete a mano)
    p_alumno = doc.add_paragraph()
    p_alumno.add_run("Nombre del alumno: ").bold = True
    p_alumno.add_run("______________________________")

    # Fecha
    p_fecha = doc.add_paragraph()
    p_fecha.add_run("Fecha: ").bold = True
    p_fecha.add_run(fecha)

    doc.add_paragraph()  # espacio antes del contenido


# ── planilla ──────────────────────────────────────────────────────────────────

@router.post("/planilla")
async def generar_planilla(tema: str, id_docente: str):
    try:
        prompt = f"Generá un cronograma de clases para el tema {tema}"
        datos_json = await RAGOrchestrator.get_context_and_generate(prompt, SYSTEM_PROMPT_XLSX)
        archivo_binario = FileEngine.create_xlsx(datos_json)
        return await process_and_upload(archivo_binario, f"Plan_{tema}", tema, "xlsx", id_docente, "RAG_XLSX")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ── apunte ────────────────────────────────────────────────────────────────────

@router.post("/apunte")
async def generar_apunte(
    tema: str = Form(...),
    id_docente: str = Form(...),
    file: Optional[UploadFile] = File(None),
    id_escuela: Optional[str] = Form(None),
    id_curso: Optional[str] = Form(None),
    fecha: Optional[str] = Form(None),
):
    try:
        nombre_escuela, nombre_materia, division, contenido_minimo, bibliografia = \
            _datos_escuela_materia(id_escuela, id_curso)

        SYSTEM_PROMPT_APUNTE = """
        Sos un asistente pedagógico experto en sintetizar material académico.
        A partir del contenido proporcionado, generá un APUNTE de estudio claro,
        enfocado en el tema indicado, listo para que un alumno lo use para repasar.

        Devolvé EXCLUSIVAMENTE un JSON válido con esta estructura:
        {
          "titulo": "Título del apunte",
          "introduccion": "Párrafo breve que contextualiza el tema",
          "secciones": [
            {
              "subtitulo": "Nombre de la sección",
              "contenido": "Explicación clara y desarrollada del concepto",
              "puntos_clave": ["idea 1", "idea 2", "idea 3"]
            }
          ],
          "glosario": [
            {"termino": "Concepto", "definicion": "Definición breve"}
          ],
          "conclusion": "Cierre con las ideas más importantes para recordar"
        }
        Mínimo 4 secciones y 5 términos en el glosario.
        """

        prompt_base = (
            f"Generá un APUNTE / RESUMEN de estudio sobre el tema '{tema}' "
            f"para la materia '{nombre_materia}' de '{nombre_escuela}'."
        )

        if file is not None:
            pdf_content = await file.read()
            datos_json = await RAGOrchestrator.get_context_from_file_and_generate(
                pdf_content, prompt_base, SYSTEM_PROMPT_APUNTE
            )
        else:
            ctx = _contexto_biblio(contenido_minimo, bibliografia)
            datos_json = await RAGOrchestrator.get_context_and_generate(
                f"{prompt_base}\n\n{ctx}", SYSTEM_PROMPT_APUNTE
            )

        from docx import Document

        doc = Document()
        _encabezado_documento(doc, "Apunte", nombre_materia, tema, nombre_escuela, fecha_str=fecha)

        if datos_json.get("introduccion"):
            doc.add_heading("Introducción", level=1)
            doc.add_paragraph(datos_json["introduccion"])

        for seccion in datos_json.get("secciones", []):
            doc.add_heading(seccion.get("subtitulo", "Sección"), level=1)
            if seccion.get("contenido"):
                doc.add_paragraph(seccion["contenido"])
            for punto in seccion.get("puntos_clave", []):
                doc.add_paragraph(punto, style="List Bullet")

        if datos_json.get("glosario"):
            doc.add_heading("Glosario", level=1)
            for item in datos_json["glosario"]:
                p = doc.add_paragraph()
                p.add_run(f"{item.get('termino', '')}: ").bold = True
                p.add_run(item.get("definicion", ""))

        if datos_json.get("conclusion"):
            doc.add_heading("Conclusión", level=1)
            doc.add_paragraph(datos_json["conclusion"])

        buffer = BytesIO()
        doc.save(buffer)
        buffer.seek(0)

        nombre_archivo = f"Apunte_{nombre_escuela}_{nombre_materia}".replace(" ", "_")
        return await process_and_upload(
            buffer.getvalue(), nombre_archivo, tema, "docx", id_docente, "RAG_APUNTE"
        )

    except Exception as e:
        print(f"DEBUG ERROR APUNTE: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


# ── preguntas guía ────────────────────────────────────────────────────────────

@router.post("/preguntas")
async def generar_preguntas(
    tema: str = Form(...),
    id_docente: str = Form(...),
    nombre_guia: str = Form(...),
    numero_preguntas: int = Form(10),
    file: Optional[UploadFile] = File(None),
    id_escuela: Optional[str] = Form(None),
    id_curso: Optional[str] = Form(None),
    fecha: Optional[str] = Form(None),
):
    from docx import Document

    try:
        if numero_preguntas < 1 or numero_preguntas > 25:
            raise HTTPException(status_code=400, detail="La cantidad de preguntas debe estar entre 1 y 25.")

        nombre_escuela, nombre_materia, division, contenido_minimo, bibliografia = \
            _datos_escuela_materia(id_escuela, id_curso)

        SYSTEM_PROMPT_PREGUNTAS = f"""
        Sos un asistente pedagógico experto en generar guías de preguntas
        a partir de un material de lectura.

        Generá EXACTAMENTE {numero_preguntas} preguntas guía sobre el tema indicado.
        Combiná preguntas de comprensión literal, de análisis y de reflexión.
        No inventes contenido que no esté en el material provisto.

        IMPORTANTE: Devolvé EXCLUSIVAMENTE un JSON válido. Estructura EXACTA:
        {{
          "titulo": "Título de la guía",
          "introduccion": "Indicaciones breves para el alumno",
          "preguntas": [
            {{"numero": 1, "pregunta": "Texto de la pregunta", "respuesta_sugerida": "Respuesta orientativa"}}
          ]
        }}
        El array 'preguntas' debe tener EXACTAMENTE {numero_preguntas} elementos.
        """

        prompt_base = (
            f"Generá una guía de {numero_preguntas} preguntas sobre el tema '{tema}' "
            f"para la materia '{nombre_materia}' de '{nombre_escuela}'."
        )

        if file is not None:
            pdf_content = await file.read()
            if not pdf_content:
                raise HTTPException(status_code=400, detail="El PDF está vacío.")
            datos_json = await RAGOrchestrator.get_context_from_file_and_generate(
                pdf_content, prompt_base, SYSTEM_PROMPT_PREGUNTAS
            )
        else:
            ctx = _contexto_biblio(contenido_minimo, bibliografia)
            datos_json = await RAGOrchestrator.get_context_and_generate(
                f"{prompt_base}\n\n{ctx}", SYSTEM_PROMPT_PREGUNTAS
            )

        # Normalización defensiva
        if isinstance(datos_json, str):
            txt = datos_json.strip()
            m = re.search(r"\{.*\}", txt, re.DOTALL)
            if m:
                try:
                    datos_json = json.loads(m.group(0))
                except Exception:
                    datos_json = {}
        if not isinstance(datos_json, dict):
            datos_json = {}

        preguntas    = datos_json.get("preguntas") or datos_json.get("questions") or datos_json.get("items") or []
        titulo       = datos_json.get("titulo") or datos_json.get("title") or nombre_guia
        introduccion = datos_json.get("introduccion") or datos_json.get("introduction") or ""

        if not preguntas:
            raise HTTPException(
                status_code=500,
                detail=f"El LLM no devolvió preguntas. Respuesta cruda: {str(datos_json)[:300]}"
            )

        doc = Document()
        _encabezado_documento(doc, "Guía de Preguntas", nombre_materia, tema, nombre_escuela, fecha_str=fecha)

        if introduccion:
            doc.add_heading("Indicaciones", level=1)
            doc.add_paragraph(introduccion)

        doc.add_heading("Preguntas", level=1)
        for i, item in enumerate(preguntas, start=1):
            if isinstance(item, dict):
                texto  = item.get("pregunta") or item.get("question") or item.get("texto") or ""
                numero = item.get("numero", i)
            else:
                texto, numero = str(item), i
            p = doc.add_paragraph()
            p.add_run(f"{numero}. ").bold = True
            p.add_run(texto)

        doc.add_page_break()
        doc.add_heading("Respuestas sugeridas (uso docente)", level=1)
        for i, item in enumerate(preguntas, start=1):
            if isinstance(item, dict):
                texto  = item.get("pregunta") or item.get("question") or ""
                resp   = item.get("respuesta_sugerida") or item.get("respuesta") or item.get("answer") or ""
                numero = item.get("numero", i)
            else:
                texto, resp, numero = str(item), "", i
            p = doc.add_paragraph()
            p.add_run(f"{numero}. {texto}\n").bold = True
            p.add_run(resp)

        buffer = BytesIO()
        doc.save(buffer)
        buffer.seek(0)

        nombre_archivo = f"Guia_{nombre_escuela}_{nombre_materia}_{nombre_guia}".replace(" ", "_")
        return await process_and_upload(
            buffer.getvalue(), nombre_archivo, tema, "docx", id_docente, "RAG_PREGUNTAS"
        )

    except HTTPException:
        raise
    except Exception as e:
        print(f"DEBUG ERROR PREGUNTAS: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


# ── examen ────────────────────────────────────────────────────────────────────

@router.post("/examen")
async def generar_examen(
    id_docente: str = Form(...),
    materia: Optional[str] = Form(None),
    fecha_examen: str = Form(...),
    tipos: str = Form(...),
    file: Optional[UploadFile] = File(None),
    id_escuela: Optional[str] = Form(None),
    id_curso: Optional[str] = Form(None),
):
    from docx import Document

    try:
        try:
            tipos_dict = json.loads(tipos)
        except Exception:
            raise HTTPException(status_code=400, detail="Formato inválido en 'tipos'.")

        seleccionados = {
            k: v for k, v in tipos_dict.items()
            if v.get("activo") and int(v.get("cantidad", 0)) > 0
        }
        if not seleccionados:
            raise HTTPException(status_code=400, detail="Marcá al menos una actividad con cantidad > 0.")

        # Obtener escuela y materia desde la BD
        nombre_escuela, nombre_materia_db, division, contenido_minimo, bibliografia = \
            _datos_escuela_materia(id_escuela, id_curso)

        # El campo materia del Form tiene prioridad; si no viene, usar el de la BD
        nombre_materia = (materia or "").strip() or nombre_materia_db

        # Bibliografía de fallback cuando no hay PDF
        biblio_str = ""
        if file is None:
            biblio_str = _contexto_biblio(contenido_minimo, bibliografia)

        descripcion_tipos = "\n".join([f"- {k}: {v['cantidad']} ítems" for k, v in seleccionados.items()])

        SYSTEM_PROMPT_EXAMEN = f"""
        Sos un asistente pedagógico experto. Generá un examen escrito sobre la materia '{nombre_materia}'.
        Incluí EXACTAMENTE los siguientes tipos y cantidades:
        {descripcion_tipos}

        Devolvé EXCLUSIVAMENTE un JSON válido (sin markdown):
        {{
          "titulo": "Examen de {nombre_materia}",
          "consignas": [
            {{"tipo": "desarrollo|multiple|completar|verdadero_falso",
              "enunciado": "Texto de la consigna",
              "items": ["item1", "item2", ...]}}
          ]
        }}
        """

        if file is not None:
            prompt = f"Generá el examen de '{nombre_materia}' para la fecha {fecha_examen}, basado en el PDF adjunto."
            pdf_content = await file.read()
            if not pdf_content:
                raise HTTPException(status_code=400, detail="El PDF está vacío.")
            datos_json = await RAGOrchestrator.get_context_from_file_and_generate(
                pdf_content, prompt, SYSTEM_PROMPT_EXAMEN
            )
        else:
            prompt = (
                f"Generá el examen de '{nombre_materia}' para la fecha {fecha_examen}."
                + (f"\n\n{biblio_str}" if biblio_str else "")
            )
            datos_json = await RAGOrchestrator.get_context_and_generate(prompt, SYSTEM_PROMPT_EXAMEN)

        if isinstance(datos_json, str):
            m = re.search(r"\{.*\}", datos_json, re.DOTALL)
            if m:
                try:
                    datos_json = json.loads(m.group(0))
                except Exception:
                    datos_json = {}
        if not isinstance(datos_json, dict):
            datos_json = {}

        consignas  = datos_json.get("consignas") or []
        titulo_doc = datos_json.get("titulo") or f"Examen de {nombre_materia}"

        if not consignas:
            raise HTTPException(status_code=500, detail=f"El LLM no devolvió consignas: {str(datos_json)[:300]}")

        doc = Document()
        _encabezado_documento(
            doc, "Examen", nombre_materia, nombre_materia,
            nombre_escuela, fecha_str=fecha_examen
        )

        doc.add_heading(titulo_doc, level=1)

        for idx, c in enumerate(consignas, start=1):
            tipo      = c.get("tipo", "")
            enunciado = c.get("enunciado", "")
            items     = c.get("items", []) or []
            doc.add_paragraph(f"{idx}) [{tipo}] {enunciado}", style="Heading 3")
            for j, it in enumerate(items, start=1):
                doc.add_paragraph(f"{j}. {it}")

        buffer = BytesIO()
        doc.save(buffer)
        buffer.seek(0)

        nombre_archivo = f"Examen_{nombre_escuela}_{nombre_materia}".replace(" ", "_")
        return await process_and_upload(
            buffer.getvalue(), nombre_archivo, nombre_materia, "docx", id_docente, "RAG_EXAMEN"
        )

    except HTTPException:
        raise
    except Exception as e:
        print(f"DEBUG ERROR EXAMEN: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


# ── podcast ───────────────────────────────────────────────────────────────────

@router.post("/podcast")
async def generar_podcast(
    tema: str = Form(...),
    id_docente: str = Form(...),
    file: Optional[UploadFile] = File(None),
    id_escuela: Optional[str] = Form(None),
    id_curso: Optional[str] = Form(None),
    fecha: Optional[str] = Form(None),
):
    from docx import Document

    try:
        nombre_escuela, nombre_materia, division, contenido_minimo, bibliografia = \
            _datos_escuela_materia(id_escuela, id_curso)

        SYSTEM_PROMPT_PODCAST = """
        Sos un asistente pedagógico experto en crear guiones de podcast educativo.
        A partir del contenido proporcionado, generá un GUIÓN de podcast completo,
        entretenido y pedagógico, listo para grabar.

        Devolvé EXCLUSIVAMENTE un JSON válido con esta estructura:
        {
          "titulo": "Título del episodio",
          "duracion_estimada": "15-20 minutos",
          "segmentos": [
            {
              "tipo": "introduccion|desarrollo|reflexion|conclusion",
              "titulo": "Nombre del segmento",
              "guion": "Texto completo del segmento con indicaciones entre [corchetes]"
            }
          ]
        }
        Mínimo 4 segmentos. El guión debe ser natural, conversacional y pedagógico.
        """

        prompt_base = (
            f"Generá un guión de podcast educativo sobre el tema '{tema}' "
            f"para la materia '{nombre_materia}' de '{nombre_escuela}'."
        )

        if file is not None:
            pdf_content = await file.read()
            if not pdf_content:
                raise HTTPException(status_code=400, detail="El PDF está vacío.")
            datos_json = await RAGOrchestrator.get_context_from_file_and_generate(
                pdf_content, prompt_base, SYSTEM_PROMPT_PODCAST
            )
        else:
            ctx = _contexto_biblio(contenido_minimo, bibliografia)
            datos_json = await RAGOrchestrator.get_context_and_generate(
                f"{prompt_base}\n\n{ctx}", SYSTEM_PROMPT_PODCAST
            )

        # Normalización defensiva
        if isinstance(datos_json, str):
            m = re.search(r"\{.*\}", datos_json, re.DOTALL)
            if m:
                try:
                    datos_json = json.loads(m.group(0))
                except Exception:
                    datos_json = {}
        if not isinstance(datos_json, dict):
            datos_json = {}

        titulo_ep = datos_json.get("titulo") or f"Podcast: {tema}"
        duracion  = datos_json.get("duracion_estimada") or ""
        segmentos = datos_json.get("segmentos") or []

        if not segmentos:
            raise HTTPException(
                status_code=500,
                detail=f"El LLM no devolvió segmentos. Respuesta cruda: {str(datos_json)[:300]}"
            )

        doc = Document()
        _encabezado_documento(doc, "Podcast", nombre_materia, tema, nombre_escuela, fecha_str=fecha)

        doc.add_heading(titulo_ep, level=1)
        if duracion:
            p_dur = doc.add_paragraph()
            p_dur.add_run("Duración estimada: ").bold = True
            p_dur.add_run(duracion)
        doc.add_paragraph()

        for seg in segmentos:
            tipo_seg   = (seg.get("tipo") or "").upper()
            titulo_seg = seg.get("titulo") or tipo_seg
            guion_seg  = seg.get("guion") or ""
            doc.add_heading(f"[{tipo_seg}] {titulo_seg}", level=2)
            doc.add_paragraph(guion_seg)
            doc.add_paragraph()

        buffer = BytesIO()
        doc.save(buffer)
        buffer.seek(0)

        nombre_archivo = f"Podcast_{nombre_escuela}_{nombre_materia}".replace(" ", "_")
        return await process_and_upload(
            buffer.getvalue(), nombre_archivo, tema, "docx", id_docente, "RAG_PODCAST"
        )

    except HTTPException:
        raise
    except Exception as e:
        print(f"DEBUG ERROR PODCAST: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


# ── presentacion ──────────────────────────────────────────────────────────────

@router.post("/presentacion")
async def generar_presentacion(
    tema: str = Form(...),
    id_docente: str = Form(...),
    file: UploadFile = File(None),
):
    """Genera un PPTX a partir de un tema y, opcionalmente, un PDF de contexto."""
    try:
        contenido_pdf = ""
        if file:
            try:
                try:
                    from PyPDF2 import PdfReader
                except ImportError:
                    from pypdf import PdfReader
                pdf_bytes = await file.read()
                temp_path = f"/tmp/{uuid.uuid4()}.pdf"
                with open(temp_path, "wb") as f:
                    f.write(pdf_bytes)
                reader = PdfReader(temp_path)
                contenido_pdf = "\n".join([p.extract_text() or "" for p in reader.pages])
                os.remove(temp_path)
            except Exception as e:
                print(f"⚠️ No pude leer el PDF: {e}")
                contenido_pdf = ""

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = tema
        slide.placeholders[1].text = f"Generado por Kōkua • {datetime.now().strftime('%d/%m/%Y')}"

        if contenido_pdf:
            bloques = [contenido_pdf[i:i+500] for i in range(0, min(len(contenido_pdf), 4000), 500)]
            for i, bloque in enumerate(bloques, start=1):
                s = prs.slides.add_slide(prs.slide_layouts[1])
                s.shapes.title.text = f"{tema} — Parte {i}"
                tf = s.placeholders[1].text_frame
                tf.text = bloque[:450]
                for p in tf.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(16)
        else:
            for titulo_slide in ["Introducción", "Desarrollo", "Conclusión"]:
                s = prs.slides.add_slide(prs.slide_layouts[1])
                s.shapes.title.text = titulo_slide
                s.placeholders[1].text = f"Contenido sobre {tema}…"

        nombre_archivo = f"{tema.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        os.makedirs("archivos_generados", exist_ok=True)
        ruta_local = f"archivos_generados/{nombre_archivo}"
        prs.save(ruta_local)

        tamanio_mb = round(os.path.getsize(ruta_local) / (1024 * 1024), 3)

        url_publica = None
        try:
            with open(ruta_local, "rb") as f:
                supabase.storage.from_("archivos").upload(
                    f"{id_docente}/{nombre_archivo}",
                    f.read(),
                    file_options={"content-type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"}
                )
            url_publica = supabase.storage.from_("archivos").get_public_url(f"{id_docente}/{nombre_archivo}")
        except Exception as e:
            print(f"⚠️ No pude subir a Supabase Storage: {e}")
            url_publica = f"/{ruta_local}"

        try:
            supabase.table("archivos_generados").insert({
                "id_docente":      id_docente,
                "nombre_archivo":  nombre_archivo,
                "tipo_formato":    "pptx",
                "sub_tipo":        "presentacion",
                "tema_especifico": tema,
                "prompt_origen":   f"Generar presentación sobre: {tema}",
                "fecha_creacion":  datetime.now().isoformat(),
                "categoria_ia":    "presentacion",
                "url_descarga":    url_publica,
                "uso_mb":          tamanio_mb,
                "descripcion":     f"Presentación generada automáticamente sobre el tema '{tema}'",
            }).execute()
        except Exception as e:
            print(f"⚠️ No pude registrar en BD: {e}")

        return {
            "status":         "success",
            "message":        "Presentación generada con éxito",
            "download_url":   url_publica,
            "nombre_archivo": nombre_archivo,
        }

    except Exception as e:
        print(f"❌ Error generando presentación: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error al generar presentación: {str(e)}")


# ── Modelo ────────────────────────────────────────────────────────────────────

class ExportarPDFRequest(BaseModel):
    url_docx:        str   # URL pública de Supabase Storage
    nombre_archivo:  str   # Sin extensión, ej: "Apunte_Fotosintesis"
    id_docente:      str


# ── Detección automática de conversor ────────────────────────────────────────

def _detectar_conversor() -> str:
    """Devuelve: 'libreoffice' | 'docx2pdf' | 'ninguno'"""
    if shutil.which("soffice") or shutil.which("libreoffice"):
        return "libreoffice"
    try:
        import docx2pdf  # noqa: F401
        return "docx2pdf"
    except ImportError:
        pass
    return "ninguno"


# ── Conversores individuales ──────────────────────────────────────────────────

def _con_libreoffice(docx_path: Path, out_dir: Path) -> Path:
    cmd = shutil.which("soffice") or shutil.which("libreoffice")
    result = subprocess.run(
        [cmd, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(docx_path)],
        capture_output=True, text=True, timeout=60,
    )
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice: {result.stderr}")
    pdf = out_dir / (docx_path.stem + ".pdf")
    if not pdf.exists():
        raise RuntimeError("LibreOffice no generó el PDF")
    return pdf


def _con_docx2pdf(docx_path: Path, out_dir: Path) -> Path:
    from docx2pdf import convert
    pdf = out_dir / (docx_path.stem + ".pdf")
    convert(str(docx_path), str(pdf))
    if not pdf.exists():
        raise RuntimeError("docx2pdf no generó el PDF")
    return pdf


def _con_reportlab(docx_path: Path, out_dir: Path) -> Path:
    """Fallback: extrae texto del .docx y arma un PDF simple con ReportLab."""
    from docx import Document as DocxDoc
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    pdf = out_dir / (docx_path.stem + ".pdf")
    word   = DocxDoc(str(docx_path))
    styles = getSampleStyleSheet()

    s_titulo  = ParagraphStyle("t1", parent=styles["Heading1"],  fontSize=16, spaceAfter=12)
    s_heading = ParagraphStyle("t2", parent=styles["Heading2"],  fontSize=13, spaceAfter=8)
    s_normal  = ParagraphStyle("no", parent=styles["Normal"],    fontSize=11, spaceAfter=6)

    story = []
    for p in word.paragraphs:
        txt = p.text.strip()
        if not txt:
            story.append(Spacer(1, 0.3 * cm))
            continue
        nombre = (p.style.name or "").lower()
        txt_safe = txt.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        if "heading 1" in nombre or "título 1" in nombre:
            story.append(Paragraph(txt_safe, s_titulo))
        elif "heading" in nombre or "título" in nombre:
            story.append(Paragraph(txt_safe, s_heading))
        else:
            story.append(Paragraph(txt_safe, s_normal))

    SimpleDocTemplate(
        str(pdf), pagesize=A4,
        leftMargin=2.5*cm, rightMargin=2.5*cm,
        topMargin=2.5*cm,  bottomMargin=2.5*cm,
    ).build(story)
    return pdf


# ── Endpoint: exportar a PDF ──────────────────────────────────────────────────

@router.post("/exportar_pdf")
async def exportar_pdf(body: ExportarPDFRequest):
    """
    Descarga el .docx desde Supabase Storage, lo convierte a PDF
    con el mejor conversor disponible y lo devuelve como descarga directa.
    Orden: LibreOffice → docx2pdf → ReportLab (fallback garantizado).
    """
    conversor = _detectar_conversor()
    print(f"🔧 Conversor PDF detectado: {conversor}")

    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir   = Path(tmp)
        safe_name = body.nombre_archivo.replace(" ", "_").replace("/", "-")
        docx_path = tmp_dir / f"{safe_name}.docx"

        # 1. Descargar .docx desde Supabase Storage
        try:
            async with httpx.AsyncClient(timeout=30) as client:
                resp = await client.get(body.url_docx)
                resp.raise_for_status()
                docx_path.write_bytes(resp.content)
        except Exception as e:
            raise HTTPException(502, detail=f"No se pudo descargar el archivo: {e}")

        if docx_path.stat().st_size == 0:
            raise HTTPException(400, detail="El archivo descargado está vacío.")

        # 2. Convertir
        try:
            if conversor == "libreoffice":
                pdf_path = _con_libreoffice(docx_path, tmp_dir)
            elif conversor == "docx2pdf":
                pdf_path = _con_docx2pdf(docx_path, tmp_dir)
            else:
                pdf_path = _con_reportlab(docx_path, tmp_dir)
        except Exception as e:
            print(f"⚠️ {conversor} falló ({e}), intentando ReportLab…")
            try:
                pdf_path = _con_reportlab(docx_path, tmp_dir)
            except Exception as e2:
                raise HTTPException(500, detail=f"Conversión fallida: {e} | Fallback: {e2}")

        # 3. Leer bytes antes de que TemporaryDirectory se cierre
        pdf_bytes = pdf_path.read_bytes()

    return StreamingResponse(
        BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={
            "Content-Disposition": f'attachment; filename="{safe_name}.pdf"',
            "Content-Length":      str(len(pdf_bytes)),
        },
    )


# ── Endpoint: diagnóstico del conversor ──────────────────────────────────────

@router.get("/exportar_pdf/info")
async def info_conversor():
    """
    Devuelve qué conversor detectó el servidor.
    Llamalo después de deployar para verificar la calidad de conversión.
    """
    conversor = _detectar_conversor()
    msgs = {
        "libreoffice": "✅ Alta calidad — preserva tablas, fuentes y estilos",
        "docx2pdf":    "✅ Alta calidad — preserva formato completo",
        "ninguno":     "⚠️ Solo texto plano (ReportLab). Instalá LibreOffice: sudo apt install libreoffice",
    }
    return {
        "conversor": conversor,
        "detalle":   msgs.get(conversor, "desconocido"),
    }