import io
import pandas as pd
from pptx import Presentation
from docx import Document
from docx.shared import Pt
from app.core.database import supabase

class FileEngine:
    # --- MOTORES DE GENERACIÓN LOCAL ---

    @staticmethod
    def create_pptx(data: dict) -> io.BytesIO:
        """
        Genera un PowerPoint con carátula personalizada (RF07).
        Espera un diccionario con 'titulo_presentacion', 'colegio', 'docente', 'materia' y 'slides'.
        """
        prs = Presentation()
        
        # 1. DIAPOSITIVA DE CARÁTULA (Personalización Automática)
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        # Título principal de la clase
        title.text = data.get("titulo_presentacion", data.get("titulo", "Clase Kōkua"))
        
        # Subtítulo con datos del docente (Magia del RF07)
        info_docente = (
            f"Institución: {data.get('colegio', 'UNDeC')}\n"
            f"Docente: {data.get('docente', 'Prof. Invitado')}\n"
            f"Materia: {data.get('materia', 'General')}"
        )
        subtitle.text = info_docente

        # 2. DIAPOSITIVAS DE CONTENIDO (Bucle sobre la lista de slides generada por IA)
        for s_data in data.get("slides", []):
            slide_layout = prs.slide_layouts[1] # Título y Contenido
            slide = prs.slides.add_slide(slide_layout)
            
            slide.shapes.title.text = s_data.get("subtitulo", s_data.get("titulo_slide", ""))
            
            tf = slide.placeholders[1].text_frame
            tf.word_wrap = True
            
            for item in s_data.get("contenido", []):
                p = tf.add_paragraph()
                p.text = str(item)
                p.level = 0 # Primer nivel de viñeta

        # Guardar en memoria
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io

    @staticmethod
    def create_docx(data: dict) -> io.BytesIO:
        doc = Document()
        doc.add_heading(data.get("titulo", "Guía de Estudio"), 0)
        for seccion in data.get("secciones", []):
            doc.add_heading(seccion.get("encabezado", ""), level=1)
            doc.add_paragraph(seccion.get("texto", ""))
        
        doc_io = io.BytesIO()
        doc.save(doc_io)
        doc_io.seek(0)
        return doc_io

    @staticmethod
    def create_xlsx(data: list) -> io.BytesIO:
        df = pd.DataFrame(data)
        xlsx_io = io.BytesIO()
        with pd.ExcelWriter(xlsx_io, engine='openpyxl') as writer:
            df.to_excel(writer, index=False, sheet_name='Planificacion_Docente')
        xlsx_io.seek(0)
        return xlsx_io

    @staticmethod
    def create_evaluacion_docx(data: dict) -> io.BytesIO:
        doc = Document()
        doc.add_heading(data.get("titulo", "Evaluación"), 0)
        for sec in data.get("secciones", []):
            doc.add_heading(sec.get("instruccion", ""), level=2)
            if sec.get("tipo") == "completar":
                for item in sec.get("items", []):
                    doc.add_paragraph(f"• {item} ____________________")
            elif sec.get("tipo") == "multiple_choice":
                for p in sec.get("preguntas", []):
                    doc.add_paragraph(p["pregunta"], style='List Bullet')
                    for opt in p["opciones"]:
                        doc.add_paragraph(f"    [  ] {opt}")
        
        doc_io = io.BytesIO()
        doc.save(doc_io)
        doc_io.seek(0)
        return doc_io

    # --- FUNCIONES DE NUBE (STORAGE) ---

    @staticmethod
    def upload_to_supabase(file_io, file_name: str, id_docente: str):
        """Sube el archivo y devuelve una URL firmada de descarga."""
        try:
            contenido = file_io.getvalue() if hasattr(file_io, "getvalue") else file_io
            path_on_storage = f"{id_docente}/{file_name}"

            supabase.storage.from_("documentos_docentes").upload(
                path=path_on_storage,
                file=contenido,
                file_options={"content-type": "application/octet-stream", "x-upsert": "true"}
            )

            # 🔑 Generar URL firmada válida por 7 días (604800 segundos)
            signed = supabase.storage.from_("documentos_docentes").create_signed_url(
                path_on_storage, 60 * 60 * 24 * 7
            )
            signed_url = signed.get("signedURL") or signed.get("signedUrl") or signed.get("signed_url")

            return {"path": path_on_storage, "url": signed_url}
        except Exception as e:
            print(f"Error subiendo a la nube: {e}")
            return None

    @staticmethod
    def get_storage_usage(id_docente: str):
        """Gestión de Límites y Termómetro de Tiza"""
        try:
            res = supabase.storage.from_("documentos_docentes").list(id_docente)
            if not res:
                return {"usado_mb": 0, "disponible_mb": 50, "porcentaje_uso": 0}

            total_bytes = sum(item['metadata'].get('size', 0) for item in res if 'metadata' in item)
            total_mb = total_bytes / (1024 * 1024)
            
            limite_usuario_mb = 50 
            disponible_mb = max(0, limite_usuario_mb - total_mb)
            
            return {
                "usado_mb": round(total_mb, 2),
                "disponible_mb": round(disponible_mb, 2),
                "porcentaje_uso": round((total_mb / limite_usuario_mb) * 100, 1)
            }
        except Exception:
            return {"usado_mb": 0, "disponible_mb": 50, "porcentaje_uso": 0}